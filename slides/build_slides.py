"""Build Pass 2 slides using Peter's PowerPoint template."""

from pptx import Presentation
from pptx.util import Pt

# Use Peter's PPT as template (preserves theme, layouts, master slides)
prs = Presentation("Koda som ett Ess 1.pptx")

# Remove Peter's slides but keep layouts
for sldId in list(prs.slides._sldIdLst):
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)


def get_layout(name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def add_bold_body_text(tf, items: list[tuple[str, str]]):
    """Add items as bold heading + normal text pairs."""
    for i, (heading, body) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = heading
        run.font.bold = True
        run.font.size = Pt(14)

        run2 = p.add_run()
        run2.text = "\n" + body
        run2.font.size = Pt(14)

        if i < len(items) - 1:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(4)


def add_bullet_list(tf, heading: str, items: list[str], *, heading_size=Pt(16), item_size=Pt(13)):
    """Add a bold heading followed by bulleted items."""
    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = heading
    run.font.bold = True
    run.font.size = heading_size
    for item in items:
        p = tf.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = item
        run.font.size = item_size


# =====================================================================
# Slide 1: Title
# =====================================================================
slide = prs.slides.add_slide(get_layout("Start och slutsida"))
title = slide.placeholders[0]
p = title.text_frame.paragraphs[0]
p.clear()
run = p.add_run()
run.text = "Koda som ett Ess"
p2 = title.text_frame.add_paragraph()
run2 = p2.add_run()
run2.text = "Pass 2: AI som medarbetare"
run2.font.size = Pt(24)
p3 = title.text_frame.add_paragraph()
run3 = p3.add_run()
run3.text = "Peter Hansson, Novatrox"
run3.font.size = Pt(24)
p4 = title.text_frame.add_paragraph()
run4 = p4.add_run()
run4.text = "Anders Viljosson"
run4.font.size = Pt(24)


# =====================================================================
# Slide 2: Övning 2a — Kontext: jämförelse
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 2 spalt text beige bkgr"))
slide.placeholders[0].text = "Övning 2: Kontextfil är central"

left = slide.placeholders[1]
left.text_frame.clear()
p = left.text_frame.paragraphs[0]
run = p.add_run()
run.text = "UTAN kontext"
run.font.bold = True
run.font.size = Pt(16)

for item in [
    "Engelska texter",
    "Inga type hints",
    "Random struktur",
    "pip install X",
    "Ingen felhantering",
]:
    p = left.text_frame.add_paragraph()
    p.level = 1
    run = p.add_run()
    run.text = "✗  " + item
    run.font.size = Pt(13)

p = left.text_frame.add_paragraph()
p.space_before = Pt(16)
run = p.add_run()
run.text = "MED CLAUDE.md (15 rader)"
run.font.bold = True
run.font.size = Pt(16)

for item in [
    "Svenska hjälptexter",
    "Type hints överallt",
    "Tydliga exit-koder",
    "Bara standardbiblioteket",
    "Docstrings, argparse",
    'if __name__ == "__main__"',
]:
    p = left.text_frame.add_paragraph()
    p.level = 1
    run = p.add_run()
    run.text = "✓  " + item
    run.font.size = Pt(13)

right = slide.placeholders[10]
right.text_frame.clear()
add_bold_body_text(right.text_frame, [
    ("Samma prompt — helt olika resultat",
     "Kontextfilen styr AI:ns beteende utan att du ändrar prompten."),
    ("Teamet delar konventioner via git",
     "Alla får samma AI-beteende, oavsett om de kör Gemini, Claude eller Cursor."),
    ("Lägg till en ny regel — AI följer automatiskt",
     "Steg 4: lägg till --json/--text i kontextfilen. AI följer det utan att du nämner det."),
])


# =====================================================================
# Slide 3: Övning 2b — Kontext: vinster
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 1 punktspalt beige bkgr"))
slide.placeholders[0].text = "Kontextfil: varför det fungerar"

body = slide.placeholders[1]
body.text_frame.clear()
add_bold_body_text(body.text_frame, [
    ("15 rader kontext > 150 rader prompt-engineering",
     "Kontexten är bestående — prompten är tillfällig. Skriv en gång, används hundratals gånger."),
    ("Onboarding för AI = onboarding för nya utvecklare",
     "Samma .md-fil som hjälper AI:n hjälper också nya teammedlemmar förstå konventionerna."),
    ("Fungerar med alla verktyg",
     "CLAUDE.md, GEMINI.md, .cursor/rules — samma innehåll, olika filnamn. Alla i teamet delar via git."),
    ("Investera 10 min, spara timmar",
     "Varje framtida interaktion med AI:n följer teamets konventioner — utan att du behöver upprepa dem."),
])


# =====================================================================
# Slide 4: Övning 3a — TDD: arbetsdelning
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 2 spalt text vit bkgr"))
slide.placeholders[0].text = "Övning 3: TDD med AI"

left = slide.placeholders[1]
left.text_frame.clear()
p = left.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Ny arbetsdelning"
run.font.bold = True
run.font.size = Pt(16)

for label, desc in [
    ("Traditionellt:", "Människa skriver spec → implementerar → testar → fixar"),
    ("Med AI:", "Människa skriver TESTER → AI implementerar → AI kör + fixar → grönt"),
]:
    p = left.text_frame.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = label
    run.font.bold = True
    run.font.size = Pt(13)
    p = left.text_frame.add_paragraph()
    run = p.add_run()
    run.text = desc
    run.font.size = Pt(13)

add_bullet_list(left.text_frame, "Iterativ loop", [
    "26 tester → AI läser → implementerar",
    "Kör pytest → fixar fel → kör igen",
    "1–3 iterationer till grönt",
])

right = slide.placeholders[10]
right.text_frame.clear()
add_bold_body_text(right.text_frame, [
    ("Tester är den mest exakta prompten",
     "Ingen tvetydighet — körbar, verifierbar, ingen \"lost in translation\"."),
    ("Du fokuserar på VAD, AI:n på HUR",
     "Du specificerar affärsregler som tester. AI väljer implementation."),
    ("AI itererar snabbare än du",
     "Implementera → test → fix → test. Agentisk loop tills allt är grönt."),
])


# =====================================================================
# Slide 5: Övning 3b — TDD: vinster och demo
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 1 punktspalt vit bkgr"))
slide.placeholders[0].text = "TDD med AI: vinster"

body = slide.placeholders[1]
body.text_frame.clear()
add_bold_body_text(body.text_frame, [
    ("Fungerar för logik OCH GUI",
     "26 backend-tester + 13 tkinter UI-tester. Samma princip, helt olika lager."),
    ("Ny affärsregel? Skriv ETT test",
     "test_weekend_surcharge → AI implementerar helgpåslag på 30 sekunder."),
    ("100% testtäckning från dag 1",
     "Testerna existerar innan koden — du börjar aldrig med teknisk skuld."),
    ("Perfekt arbetsdelning",
     "Människa specificerar, AI implementerar och itererar. Du behåller kontrollen via testerna."),
])


# =====================================================================
# Slide 6: Övning 4a — AI-review: arbetsflöde
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 2 spalt text beige bkgr"))
slide.placeholders[0].text = "Övning 4: AI-review mot kodstandarden"

left = slide.placeholders[1]
left.text_frame.clear()
p = left.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Professionellt arbetsflöde"
run.font.bold = True
run.font.size = Pt(16)

for step in [
    "1. Fork + branch",
    "2. Koda med AI (Gemini/Claude)",
    "3. Commit + push + öppna PR",
    "4. AI-review via GitHub Action",
    "5. Fixa avvikelser → ny review",
    "6. Merge när godkänt",
]:
    p = left.text_frame.add_paragraph()
    p.level = 1
    run = p.add_run()
    run.text = step
    run.font.size = Pt(13)

add_bullet_list(left.text_frame, "Tvåsidigt AI-skydd", [
    "Skrivande AI följer kontextfilen (övning 2)",
    "Tester fångar logikfel (övning 3)",
    "Granskande AI fångar stilfel mot kodstandard.md",
    "Merge blockeras vid avvikelser",
])

right = slide.placeholders[10]
right.text_frame.clear()
add_bold_body_text(right.text_frame, [
    ("Automatisk kvalitetsgrind på varje PR",
     "Kodstandarden genomdrivs konsekvent — ingen \"det där låter vi gå\"."),
    ("Olika AI vid skrivning vs granskning",
     "Fräscha ögon fångar fler problem. AI:n som skrev koden hittar sällan sina egna misstag."),
    ("Speglar riktigt utvecklararbetsflöde",
     "Fork → branch → PR → review. Samma process som på jobbet, med AI i varje steg."),
])


# =====================================================================
# Slide 7: Övning 4b — AI-review: vinster
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 1 punktspalt beige bkgr"))
slide.placeholders[0].text = "AI-review: varför det fungerar"

body = slide.placeholders[1]
body.text_frame.clear()
add_bold_body_text(body.text_frame, [
    ("Du är granskare, inte skribent",
     "Högre nivå — du styr arkitektur och kvalitet istället för att skriva varje rad."),
    ("AI ersätter inte mänsklig granskning",
     "Men den fångar det du missar. Precis som CI-tester är en extra säkerhetsnivå."),
    ("Kontext styr även granskningen",
     "Utan kodstandard.md ger AI:n generisk feedback. Med den pekar den ut specifika avvikelser mot ert regelverk."),
    ("Övning 4 knyter ihop allt",
     "Kontext (övning 2) + TDD (övning 3) + agentiskt arbete + automatisk granskning = komplett arbetsflöde."),
])


# =====================================================================
# Slide 8: Sammanfattning
# =====================================================================
slide = prs.slides.add_slide(get_layout("Rubrik 1 punktspalt vit bkgr"))
slide.placeholders[0].text = "Tre principer att ta med"

body = slide.placeholders[1]
body.text_frame.clear()
add_bold_body_text(body.text_frame, [
    ("1. Kontext först (övning 2)",
     "15 rader i en .md-fil > clevera prompts. Hela teamet får konsekvent AI-beteende."),
    ("2. Tester som spec (övning 3)",
     "Skriv vad du vill ha, låt AI skriva hur. Människa specificerar, AI implementerar och itererar."),
    ("3. AI som kollega, inte verktyg (övning 4)",
     "Ge uppdraget, låt agenten jobba. Du granskar — högre nivå, bättre resultat."),
    ("Gemensam nämnare",
     "Du höjer dig från SKRIBENT till ARKITEKT och GRANSKARE."),
])


# =====================================================================
# Slide 9: End slide
# =====================================================================
slide = prs.slides.add_slide(get_layout("Start och slutsida"))
title = slide.placeholders[0]
title.text_frame.clear()
p = title.text_frame.paragraphs[0]
run = p.add_run()
run.text = "peter.hansson@novatrox.se"
run.font.size = Pt(24)
p2 = title.text_frame.add_paragraph()
run2 = p2.add_run()
run2.text = "anders@viljo.se"
run2.font.size = Pt(24)


prs.save("Koda som ett Ess 2.pptx")
print("Saved: Koda som ett Ess 2.pptx")
